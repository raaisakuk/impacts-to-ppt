from __future__ import print_function
import os

import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt

path_dir = os.path.join(os.path.dirname(__file__))

def df_to_table(slide, df, case_name=None, left=Inches(5), top=Inches(1), width=Inches(7), height=Inches(5), colnames=None):
    '''Converts a Pandas DataFrame to a PowerPoint table on the given
    Slide of a PowerPoint presentation.
    The table is a standard Powerpoint table, and can easily be modified with the Powerpoint tools,
    for example: resizing columns, changing formatting etc.
    Arguments:
     - slide: slide object from the python-pptx library containing the slide on which you want the table to appear
     - df: Pandas DataFrame with the data
    Optional arguments:
     - colnames
     https://github.com/robintw/PandasToPowerpoint/blob/master/PandasToPowerpoint.py

    Copyright (c) 2015, Chris Moffitt
    All rights reserved.

    Redistribution and use in source and binary forms, with or without
    modification, are permitted provided that the following conditions are met:

    * Redistributions of source code must retain the above copyright notice, this
      list of conditions and the following disclaimer.

    * Redistributions in binary form must reproduce the above copyright notice,
      this list of conditions and the following disclaimer in the documentation
      and/or other materials provided with the distribution.

    * Neither the name of pbpython nor the names of its
      contributors may be used to endorse or promote products derived from
      this software without specific prior written permission.

    THIS SOFTWARE IS PROVIDED BY THE COPYRIGHT HOLDERS AND CONTRIBUTORS "AS IS"
    AND ANY EXPRESS OR IMPLIED WARRANTIES, INCLUDING, BUT NOT LIMITED TO, THE
    IMPLIED WARRANTIES OF MERCHANTABILITY AND FITNESS FOR A PARTICULAR PURPOSE ARE
    DISCLAIMED. IN NO EVENT SHALL THE COPYRIGHT HOLDER OR CONTRIBUTORS BE LIABLE
    FOR ANY DIRECT, INDIRECT, INCIDENTAL, SPECIAL, EXEMPLARY, OR CONSEQUENTIAL
    DAMAGES (INCLUDING, BUT NOT LIMITED TO, PROCUREMENT OF SUBSTITUTE GOODS OR
    SERVICES; LOSS OF USE, DATA, OR PROFITS; OR BUSINESS INTERRUPTION) HOWEVER
    CAUSED AND ON ANY THEORY OF LIABILITY, WHETHER IN CONTRACT, STRICT LIABILITY,
    OR TORT (INCLUDING NEGLIGENCE OR OTHERWISE) ARISING IN ANY WAY OUT OF THE USE
    OF THIS SOFTWARE, EVEN IF ADVISED OF
    THE POSSIBILITY OF SUCH DAMAGE.

    See http://pbpython.com/creating-powerpoint.html for details on this script
    Requires https://python-pptx.readthedocs.org/en/latest/index.html

    Example program showing how to read in Excel, process with pandas and
    output to a PowerPoint file.
    '''
    rows, cols = df.shape
    res = slide.shapes.add_table(rows + 1, cols, left, top, width, height)

    if colnames is None:
        colnames = list(df.columns)

    if case_name:
        colnames[0] = case_name
    # Insert the column names
    for col_index, col_name in enumerate(colnames):
        # Column names can be tuples
        if not isinstance(col_name, str):
            # col_name = " ".join(col_name)
            col_name = str(col_name)
        res.table.cell(0, col_index).text = col_name

    m = df.as_matrix()

    for row in range(rows):
        for col in range(cols):
            val = m[row, col]
            text = str(val)
            res.table.cell(row + 1, col).text = text
            para = res.table.cell(row + 1, col).text_frame.paragraphs[0]
            para.font.size = Pt(12)

def create_ppt(output, case_names, case_figs, case_tables, emsc_fig, emsc_parts_fig, overall_fig, overall_table):
    """ Take the input powerpoint file and use it as the template for the output
    file.
    """
    prs = Presentation(os.path.join(path_dir, 'report_out.pptx'))

    #Slide 1
    perf_summ = prs.slide_layouts[0]
    slide = prs.slides.add_slide(perf_summ)
    overall_fig.savefig('overall_fig.png')
    top = Inches(1)
    left = Inches(1)
    width = Inches(10)
    height = Inches(5)
    slide.shapes.add_picture('overall_fig.png', left, top, width, height)

    #Slide 2
    emsc = prs.slide_layouts[0]
    slide = prs.slides.add_slide(emsc)
    emsc_fig.savefig('emsc_fig.png')
    top = Inches(1)
    left = Inches(1)
    width = Inches(3)
    height = Inches(3)
    slide.shapes.add_picture('emsc_fig.png', left, top, width, height)

    emsc_parts_fig.savefig('emsc_parts_fig.png')
    top = Inches(1)
    left = Inches(5)
    width = Inches(7)
    height = Inches(5)
    slide.shapes.add_picture('emsc_parts_fig.png', left, top, width, height)

    #Slide 3
    for i in range(len(case_figs)):
        if isinstance(case_tables[i], pd.DataFrame):
            case_figs[i].savefig('case_'+str(i)+'.png')
            case_slide = prs.slide_layouts[0]
            slide = prs.slides.add_slide(case_slide)
            top = Inches(1)
            left = Inches(1)
            width = Inches(3)
            height = Inches(3)
            slide.shapes.add_picture('case_'+str(i)+'.png', left, top, width, height)
            if(len(case_tables[i])>12):
                df_to_table(slide, case_tables[i][:12], case_names[i])
                case_slide = prs.slide_layouts[0]
                slide = prs.slides.add_slide(case_slide)
                df_to_table(slide, case_tables[i][12:], case_names[i])
                top = Inches(1)
                left = Inches(1)
                width = Inches(3)
                height = Inches(3)
                slide.shapes.add_picture('case_'+str(i)+'.png', left, top, width, height)
            else:
                df_to_table(slide, case_tables[i], case_names[i])

    #Slide 4
    overall_perf = prs.slide_layouts[0]
    slide = prs.slides.add_slide(overall_perf)
    df_to_table(slide=slide, df=overall_table, left=Inches(3))

    prs.save(output)



